"""Generate presentation from report results."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

OUT = "presentation.pptx"

# ── Colour palette ────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x4F, 0x8A)   # dark blue  — headings
C_MID    = RGBColor(0x4C, 0x9B, 0xE8)   # mid blue   — accents
C_LIGHT  = RGBColor(0xEE, 0xF4, 0xFB)   # pale blue  — shaded boxes
C_RED    = RGBColor(0xC0, 0x39, 0x2B)   # red        — night regression
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)   # green      — positive delta
C_ORANGE = RGBColor(0xE8, 0xA2, 0x3C)   # orange     — night
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(BLANK)

def box(slide, x, y, w, h, text='', size=18, bold=False, italic=False,
        color=None, bg=None, align=PP_ALIGN.LEFT, wrap=True, valign=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = 'Calibri'
    run.font.color.rgb = color or C_TEXT
    if bg:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg
    return txBox

def rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.15, C_DARK)
    box(slide, 0.35, 0.08, 12.5, 0.65, title, size=28, bold=True, color=C_WHITE)
    if subtitle:
        box(slide, 0.35, 0.72, 12.5, 0.38, subtitle, size=14, color=C_MID)

def bullet_list(slide, items, x, y, w, h, size=16, indent=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        if indent:
            p.level = 1
        run = p.add_run()
        run.text = ('• ' if not indent else '  – ') + item
        run.font.size = Pt(size)
        run.font.name = 'Calibri'
        run.font.color.rgb = C_TEXT
        p.space_after = Pt(6)

def add_image(slide, path, x, y, w):
    try:
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    except Exception as e:
        box(slide, x, y, w, 2, f'[Image: {path}]', size=11, italic=True, color=C_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, C_DARK)
rect(sl, 0, 2.5, 13.33, 2.7, RGBColor(0x0F, 0x30, 0x57))
box(sl, 0.8, 0.4, 11.7, 0.5, 'MSc Computer Vision', size=16, color=C_MID, align=PP_ALIGN.CENTER)
box(sl, 0.8, 1.0, 11.7, 1.3,
    'Two-Step Unsupervised Domain Adaptation\nfor Semantic Segmentation Under Adverse Weather',
    size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
box(sl, 0.8, 2.6, 11.7, 0.5, 'SegFormer-B2  ·  Mean Teacher  ·  GTA5 → Cityscapes → ACDC',
    size=16, color=C_MID, align=PP_ALIGN.CENTER)
box(sl, 0.8, 3.3, 11.7, 0.4, 'Gal Hanuna', size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
box(sl, 0.8, 3.75, 11.7, 0.4, 'May 2026', size=15, color=C_MID, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Motivation
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Motivation', 'Why adverse-weather UDA matters')
box(sl, 0.4, 1.25, 7.5, 0.4, 'The Problem', size=17, bold=True, color=C_DARK)
bullet_list(sl, [
    'Models trained on clean/synthetic data fail under fog, rain, snow, and night',
    'Collecting labeled data for every condition and city is prohibitively expensive',
    'Deployment in new regions or weather requires re-annotation — not scalable',
], 0.4, 1.65, 7.5, 2.2, size=15)
box(sl, 0.4, 3.85, 7.5, 0.4, 'The Opportunity', size=17, bold=True, color=C_DARK)
bullet_list(sl, [
    'Large synthetic datasets (GTA5) are free to generate at scale',
    'Real clean-weather data (Cityscapes) is already labeled',
    'UDA transfers knowledge without any target-domain annotations',
], 0.4, 4.25, 7.5, 2.0, size=15)
rect(sl, 8.3, 1.2, 4.6, 5.4, C_LIGHT, C_MID)
box(sl, 8.5, 1.35, 4.2, 0.4, 'Domain Shift', size=15, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
for label, y_ in [('GTA5\n(Synthetic, clean)', 1.85), ('Cityscapes\n(Real, clean)', 3.05), ('ACDC\n(Real, adverse)', 4.25)]:
    rect(sl, 8.7, y_, 3.8, 0.85, C_DARK)
    box(sl, 8.7, y_, 3.8, 0.85, label, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if y_ < 4.25:
        box(sl, 10.1, y_+0.85, 1.0, 0.2, '↓', size=16, bold=True, color=C_MID, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Research Question & Approach
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Research Question & Two-Step Pipeline')
rect(sl, 0.4, 1.2, 12.5, 1.0, C_LIGHT, C_MID)
box(sl, 0.6, 1.3, 12.1, 0.8,
    '"Does a second adaptation step from Cityscapes to ACDC improve segmentation on ACDC,\nand which weather conditions benefit most or remain hardest?"',
    size=15, italic=True, color=C_DARK, align=PP_ALIGN.CENTER)
box(sl, 0.4, 2.4, 12.5, 0.35, 'The Two-Step Pipeline', size=17, bold=True, color=C_DARK)
for i, (label, sub, x_) in enumerate([
    ('GTA5', '24,966 labeled\nsynthetic pairs', 0.5),
    ('Step 1\nGTA5 → CS', 'Pseudo-label UDA\n20K→26.8K iters', 3.8),
    ('Cityscapes', '2,975 labeled\nreal images', 6.9),
    ('Step 2\nCS → ACDC', 'Pseudo-label UDA\n10K iters', 9.8),
    ('ACDC', '400 unlabeled\nper condition', 12.2),
]):
    if i % 2 == 0:
        rect(sl, x_, 2.85, 2.4, 1.5, C_DARK)
        box(sl, x_, 2.85, 2.4, 0.9, label, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        box(sl, x_, 3.6, 2.4, 0.75, sub, size=11, color=C_MID, align=PP_ALIGN.CENTER)
    else:
        rect(sl, x_, 3.1, 2.5, 0.9, C_MID)
        box(sl, x_, 3.1, 2.5, 0.9, label, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
box(sl, 0.4, 4.55, 12.5, 0.35, 'Key design insight:', size=14, bold=True, color=C_DARK)
box(sl, 0.4, 4.9, 12.5, 1.3,
    'The GTA5→ACDC gap is too large for one step. Factoring it into two smaller hops lets each step exploit '
    'a clear intermediate signal. Step 2 also uses Cityscapes ground-truth labels as a supervised source — '
    'free supervision that a direct single-step approach would lack.',
    size=13, color=C_TEXT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Method
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Method', 'SegFormer-B2 + Mean Teacher pseudo-labeling')
box(sl, 0.4, 1.2, 5.8, 0.35, 'SegFormer-B2 Backbone', size=16, bold=True, color=C_DARK)
bullet_list(sl, [
    '~25M parameters, Mix Transformer (MiT) encoder',
    'No positional encoding → robust to geometric distortions',
    'Multi-scale features: strides {4,8,16,32}',
    'Fits 20 GB VRAM at batch size 2 per domain',
], 0.4, 1.6, 5.8, 2.5, size=14)
box(sl, 0.4, 4.1, 5.8, 0.35, 'Training', size=16, bold=True, color=C_DARK)
bullet_list(sl, [
    'AdamW, lr=6×10⁻⁵, polynomial decay, bfloat16',
    'Pseudo-label threshold τ=0.9 (both steps)',
    'EMA decay 0.9999',
], 0.4, 4.5, 5.8, 1.8, size=14)
rect(sl, 6.6, 1.2, 6.3, 5.6, C_LIGHT, C_MID)
box(sl, 6.8, 1.3, 5.9, 0.4, 'Mean Teacher Framework', size=15, bold=True, color=C_DARK)
for label, desc, y_ in [
    ('Labeled Source', 'GTA5 (Step 1) or Cityscapes (Step 2)', 1.85),
    ('Student', 'Trained on source labels + target pseudo-labels', 2.7),
    ('Teacher (EMA)', 'Exponential moving avg of student weights', 3.55),
    ('Pseudo-labels', 'Teacher predictions on unlabeled target\n(confidence ≥ τ = 0.9 only)', 4.4),
]:
    rect(sl, 6.8, y_, 5.9, 0.65 if '\n' not in desc else 0.85, C_DARK)
    box(sl, 6.85, y_, 2.2, 0.65, label, size=12, bold=True, color=C_WHITE)
    box(sl, 9.1, y_, 3.5, 0.65, desc, size=11, color=C_MID)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Results Overview
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Results', 'Two runs — both confirm Step 2 helps')
headers = ['', 'Overall', 'Fog', 'Rain', 'Snow', 'Night']
partial_before = ['Before Step 2', '27.20%', '33.4%', '30.7%', '29.4%', '11.5%']
partial_after  = ['After Step 2',  '32.72%', '45.7%', '40.1%', '38.0%', '8.4%']
partial_delta  = ['Δ (partial)',   '+5.52 pp', '+12.3', '+9.4', '+8.6', '−3.1']
full_before    = ['Before Step 2', '32.07%', '40.54%', '35.89%', '33.38%', '15.02%']
full_after     = ['After Step 2',  '38.18%', '49.08%', '42.71%', '40.54%', '16.77%']
full_delta     = ['Δ (full-data)', '+6.11 pp', '+8.54', '+6.83', '+7.16', '+1.74']

box(sl, 0.4, 1.2, 5.9, 0.3, 'Partial-data run  (Step 1: 30.30% CS val)', size=13, bold=True, color=C_DARK)
col_w = [1.7, 0.85, 0.72, 0.72, 0.72, 0.72]
col_x = [0.4, 2.12, 2.99, 3.73, 4.47, 5.21]
for row_i, (row, bg) in enumerate(zip(
    [headers, partial_before, partial_after, partial_delta],
    [C_DARK, C_LIGHT, C_LIGHT, RGBColor(0xD5, 0xE8, 0xF5)]
)):
    y_ = 1.55 + row_i * 0.42
    for j, (txt, cx, cw) in enumerate(zip(row, col_x, col_w)):
        rect(sl, cx, y_, cw, 0.4, C_DARK if row_i == 0 else bg,
             C_MID if row_i > 0 else None)
        clr = C_WHITE if row_i == 0 else (C_RED if txt.startswith('−') else (C_GREEN if txt.startswith('+') else C_TEXT))
        box(sl, cx, y_, cw, 0.4, txt, size=11, bold=(row_i==0 or j==0),
            color=clr, align=PP_ALIGN.CENTER)

box(sl, 0.4, 3.45, 5.9, 0.3, 'Full-data run  (Step 1: 38.90% CS val)', size=13, bold=True, color=C_DARK)
for row_i, (row, bg) in enumerate(zip(
    [headers, full_before, full_after, full_delta],
    [C_DARK, C_LIGHT, C_LIGHT, RGBColor(0xD5, 0xE8, 0xF5)]
)):
    y_ = 3.8 + row_i * 0.42
    for j, (txt, cx, cw) in enumerate(zip(row, col_x, col_w)):
        rect(sl, cx, y_, cw, 0.4, C_DARK if row_i == 0 else bg,
             C_MID if row_i > 0 else None)
        clr = C_WHITE if row_i == 0 else (C_RED if txt.startswith('−') else (C_GREEN if txt.startswith('+') else C_TEXT))
        box(sl, cx, y_, cw, 0.4, txt, size=11, bold=(row_i==0 or j==0),
            color=clr, align=PP_ALIGN.CENTER)

box(sl, 6.6, 1.2, 6.3, 0.35, 'Key takeaways', size=15, bold=True, color=C_DARK)
bullet_list(sl, [
    'Step 2 always helps — +5.5 to +6.1 pp overall',
    'Fog benefits most in both runs',
    'Full-data run: every condition improves',
    'Night: regresses in partial, improves in full-data',
    'Full-data fog crosses 49% — approaching supervised range',
], 6.6, 1.6, 6.3, 3.0, size=14)
rect(sl, 6.6, 4.7, 6.3, 1.9, C_LIGHT, C_MID)
box(sl, 6.8, 4.8, 5.9, 0.35, 'Absolute mIoU context', size=13, bold=True, color=C_DARK)
box(sl, 6.8, 5.15, 5.9, 1.3,
    'Numbers reflect fully unsupervised adaptation — zero ACDC labels at any stage. '
    'The scientific contribution is the delta, not the absolute level. '
    'Night suppresses the overall average; fog/rain/snow alone reach 40–49%.',
    size=12, color=C_TEXT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Figure 1: Before/After bar chart
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Per-Condition Results', 'Before vs. After Step 2 — partial-data B2 pipeline')
add_image(sl, 'figures/fig_before_after_bar.png', 1.5, 1.2, 10.3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Adaptation curve (partial)
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Step 2 Adaptation Curve', 'Per-condition mIoU over training — partial-data run')
add_image(sl, 'figures/fig_step2_adaptation_curve.png', 1.0, 1.1, 11.3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Figure: Before/After bar chart (full-data)
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Per-Condition Results — Full-Data Pipeline', 'Before vs. After Step 2 — 24,966 GTA5 pairs, Step 1: 38.90% CS val')
add_image(sl, 'figures/fig_before_after_bar_full.png', 1.5, 1.2, 10.3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Adaptation curve (full-data)
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Step 2 Adaptation Curve — Full-Data Pipeline', 'All four conditions improve, including night — best checkpoint at iter 5,000')
add_image(sl, 'figures/fig_step2_adaptation_curve_full.png', 1.0, 1.1, 11.3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 (was 8) — Key Finding: Natural Experiment
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Key Finding: The Natural Experiment', 'One variable changed — night flips')
rect(sl, 0.4, 1.2, 12.5, 1.0, C_LIGHT, C_MID)
box(sl, 0.6, 1.28, 12.1, 0.85,
    'The two runs are identical except for Step 1 dataset size. '
    'Everything else — backbone, training procedure, Step 2 config — is the same.',
    size=15, italic=True, color=C_DARK, align=PP_ALIGN.CENTER)

for label, step1, night_before, night_after, delta, delta_color, x_ in [
    ('Partial-data', '30.30% CS val\n(~10K GTA5 pairs)', '11.5%', '8.4%', '−3.1 pp  REGRESSION', C_RED, 0.4),
    ('Full-data',   '38.90% CS val\n(24,966 GTA5 pairs)', '15.02%', '16.77%', '+1.74 pp  IMPROVEMENT', C_GREEN, 6.8),
]:
    rect(sl, x_, 2.4, 6.0, 4.7, C_LIGHT, C_MID)
    box(sl, x_+0.15, 2.5, 5.7, 0.4, label, size=16, bold=True, color=C_DARK)
    box(sl, x_+0.15, 2.95, 5.7, 0.5, f'Step 1: {step1}', size=13, color=C_GRAY)
    box(sl, x_+0.15, 3.5, 2.5, 0.35, 'Night Before:', size=12, bold=True, color=C_TEXT)
    box(sl, x_+2.7, 3.5, 3.0, 0.35, night_before, size=14, bold=True, color=C_ORANGE)
    box(sl, x_+0.15, 3.9, 2.5, 0.35, 'Night After:', size=12, bold=True, color=C_TEXT)
    box(sl, x_+2.7, 3.9, 3.0, 0.35, night_after, size=14, bold=True, color=C_ORANGE)
    rect(sl, x_+0.15, 4.4, 5.5, 0.55, C_DARK)
    box(sl, x_+0.15, 4.4, 5.5, 0.55, delta, size=16, bold=True, color=delta_color, align=PP_ALIGN.CENTER)

box(sl, 0.4, 7.0, 12.5, 0.4,
    '→  Step 1 quality is the causal mechanism for night adaptation, not the method itself.',
    size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Broader Applicability
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Broader Applicability', 'The principle extends beyond autonomous driving')
box(sl, 0.4, 1.2, 12.5, 0.5,
    'Two-step UDA applies whenever a large domain gap can be factored through a labeled intermediate domain:',
    size=15, color=C_TEXT)
for title, desc, x_, y_ in [
    ('Autonomous\nDriving', 'Synthetic game data\n→ Clean real scenes\n→ Adverse weather', 0.5, 1.85),
    ('Medical\nImaging',    'Synthetic scans\n→ One scanner type\n→ Noisier scanner', 3.5, 1.85),
    ('Satellite\nImagery',  'One season/region\n→ Intermediate domain\n→ Target region', 6.5, 1.85),
    ('Outdoor\nRobotics',   'Simulation\n→ Controlled env\n→ Unstructured env', 9.5, 1.85),
]:
    rect(sl, x_, y_, 2.7, 2.8, C_DARK)
    box(sl, x_, y_, 2.7, 0.75, title, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    box(sl, x_+0.1, y_+0.8, 2.5, 1.9, desc, size=12, color=C_MID, align=PP_ALIGN.CENTER)
rect(sl, 0.4, 4.85, 12.5, 1.5, C_LIGHT, C_MID)
box(sl, 0.6, 4.95, 12.1, 0.35, 'Structural requirement for two-step UDA:', size=14, bold=True, color=C_DARK)
box(sl, 0.6, 5.35, 12.1, 0.9,
    'Labeled synthetic source  ·  Accessible labeled intermediate domain  ·  Hard unlabeled final target\n'
    'Where the direct source→target gap is too large to bridge in one step.',
    size=13, color=C_TEXT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Conclusion
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, 'Conclusion')
bullet_list(sl, [
    'Step 2 consistently improves ACDC segmentation — +5.52 pp (partial) and +6.11 pp (full-data)',
    'Fog benefits most; rain and snow improve substantially; night is the hardest condition',
    'Full-data run: every condition improves, including night (+1.74 pp)',
    'Natural experiment: the two runs isolate Step 1 quality as the causal mechanism for night behavior',
    'Night regression is NOT a fundamental limitation — it is a consequence of a weak Step 1 model',
    'Sequential UDA factorizes hard domain gaps; Step 1 quality is the primary design lever',
    'Principle generalises beyond autonomous driving to any multi-hop domain adaptation setting',
], 0.5, 1.25, 12.3, 5.2, size=16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Thank you
# ─────────────────────────────────────────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, C_DARK)
box(sl, 0.8, 2.2, 11.7, 1.0, 'Thank you', size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
box(sl, 0.8, 3.4, 11.7, 0.5, 'Questions?', size=24, color=C_MID, align=PP_ALIGN.CENTER)
box(sl, 0.8, 4.4, 11.7, 0.4,
    'github.com/galbe/weather-uda-segformer',
    size=15, color=C_MID, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved {OUT}  ({prs.slides.__len__()} slides)")
